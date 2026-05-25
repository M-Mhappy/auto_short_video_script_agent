"""Export presentation steps to PowerPoint (.pptx)."""

from __future__ import annotations

import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from backend.tools.slide_layout import (
    ACCENT_COLOR, FG_COLOR, MUTED_COLOR, BG_COLOR, layout_from_step,
)


def _rgb(t: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*t)


def _set_slide_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(BG_COLOR)


def _add_textbox(
    slide,
    left, top, width, height,
    text: str,
    *,
    font_size: int = 28,
    color: tuple[int, int, int] = FG_COLOR,
    bold: bool = False,
    italic: bool = False,
    align=PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.color.rgb = _rgb(color)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "微软雅黑"


def _add_multiline_textbox(
    slide,
    left, top, width, height,
    lines: list[str],
    *,
    font_size: int = 28,
    color: tuple[int, int, int] = FG_COLOR,
    bold: bool = False,
    italic: bool = False,
) -> None:
    if not lines:
        return
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.color.rgb = _rgb(color)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = "微软雅黑"


def _populate_slide(slide, step: dict) -> None:
    _set_slide_bg(slide)
    layout = layout_from_step(step)
    pad = Inches(0.7)
    content_w = Inches(16) - pad * 2
    content_left = pad

    if layout.step_type == "hero":
        title_top = Inches(3.2)
        _add_textbox(
            slide, content_left, title_top, content_w, Inches(2.5),
            layout.title,
            font_size=54, bold=True, align=PP_ALIGN.CENTER,
        )
        if layout.subtitle:
            _add_textbox(
                slide, content_left, Inches(5.8), content_w, Inches(1.5),
                layout.subtitle,
                font_size=24, color=MUTED_COLOR, align=PP_ALIGN.CENTER,
            )
        return

    y = Inches(1.2)
    if layout.chapter_tag:
        _add_textbox(
            slide, content_left, y, content_w, Inches(0.6),
            layout.chapter_tag,
            font_size=14, color=ACCENT_COLOR,
        )
        y += Inches(0.7)

    if layout.step_type == "quote":
        bar = slide.shapes.add_shape(
            1, content_left, y, Inches(0.08), Inches(3.5),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = _rgb(ACCENT_COLOR)
        bar.line.fill.background()
        _add_multiline_textbox(
            slide, content_left + Inches(0.35), y, content_w - Inches(0.5), Inches(5),
            layout.body_lines,
            font_size=36, italic=True,
        )
        return

    if layout.step_type == "list_item":
        if layout.list_title:
            _add_textbox(
                slide, content_left, y, content_w, Inches(0.6),
                layout.list_title,
                font_size=18, color=MUTED_COLOR,
            )
            y += Inches(0.65)
        prefix = f"{layout.item_index}. " if layout.item_index else ""
        body = prefix + (layout.body_lines[0] if layout.body_lines else "")
        _add_textbox(
            slide, content_left, y, content_w, Inches(4),
            body,
            font_size=30, bold=True,
        )
        return

    _add_multiline_textbox(
        slide, content_left, y, content_w, Inches(6),
        layout.body_lines,
        font_size=24,
    )


def generate_pptx(steps: list[dict], book_info: dict, output_path: str) -> str:
    """Build a 16:9 PPTX with one slide per presentation step."""
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank = prs.slide_layouts[6]

    title = book_info.get("title", "")
    author = book_info.get("author", "")

    if title:
        cover = prs.slides.add_slide(blank)
        _set_slide_bg(cover)
        cover_text = f"《{title}》"
        if author:
            cover_text += f"\n{author}"
        _add_textbox(
            cover, Inches(1), Inches(3.5), Inches(14), Inches(2),
            cover_text,
            font_size=44, bold=True, align=PP_ALIGN.CENTER,
        )
        notes = cover.notes_slide.notes_text_frame
        notes.text = f"《{title}》演示文稿"

    for step in steps:
        slide = prs.slides.add_slide(blank)
        _populate_slide(slide, step)
        narration = str(step.get("narration", "")).strip()
        if narration:
            slide.notes_slide.notes_text_frame.text = narration

    prs.save(output_path)
    return output_path
